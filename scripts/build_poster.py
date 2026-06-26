"""Build poster3.pptx from template - replace method area with linear flow."""
import shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 1. Copy template
src = None
base = os.path.join(os.path.dirname(__file__), '..', 'docs')
for root, dirs, files in os.walk(base):
    for f in files:
        if 'oster' in f.lower() and f.endswith('.pptx') and '~' not in f:
            src = os.path.join(root, f); break
dst = os.path.join(os.path.dirname(__file__), '..', 'course-report', 'poster_final.pptx')
shutil.copy2(src, dst)
print('Copied template')

prs = Presentation(dst)
slide = prs.slides[0]

# 2. Delete ALL shapes in method area (T=2.0~10.5, L=4.8~19.2)
to_del = []
for s in slide.shapes:
    t = (s.top or 0) / 914400
    l = (s.left or 0) / 914400
    if 2.0 < t < 10.5 and 4.8 < l < 19.2:
        to_del.append(s)
for s in to_del:
    s._element.getparent().remove(s._element)
print(f'Deleted {len(to_del)} shapes in method area')

# Colors
ACCENT = RGBColor(0xFF, 0x46, 0x55)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1B, 0x4F, 0x9B)
GREEN = RGBColor(0x2D, 0x8B, 0x5F)
ORANGE = RGBColor(0xCC, 0x66, 0x00)
PURPLE = RGBColor(0x66, 0x33, 0x99)
FONT = 'Cambria'


def add_text(l, t, w, h, text, size, bold, color, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def add_box(l, t, w, h, fill, text, size=12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    return s


def add_arrow(l, t, w, h, direction='right', color=ACCENT):
    st = MSO_SHAPE.RIGHT_ARROW if direction == 'right' else MSO_SHAPE.DOWN_ARROW
    s = slide.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_pic(l, t, w, h, label):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
    s.line.color.rgb = RGBColor(0xBB, 0xBB, 0xCC)
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.color.rgb = GRAY
    r.font.name = FONT
    return s


# 3. Build linear flow: 5 steps in a row
add_text(5.0, 2.5, 13.5, 0.5, '\u7cfb\u7edf\u5168\u6d41\u7a0b', 24, True, ACCENT, PP_ALIGN.CENTER)

steps = [
    (5.3, '(i) \u56fe\u50cf\u91c7\u96c6', '\u5de5\u4e1a\u76f8\u673a/\u89c6\u9891', BLUE,
     '\u3010\u63d2\u56fe\uff1a\u91c7\u96c6\u753b\u9762\u3011'),
    (8.3, '(ii) YOLO26n\u68c0\u6d4b', '6\u7c7b\u7f3a\u9677+\u7f6e\u4fe1\u5ea6',
     RGBColor(0x00, 0x60, 0xA0), '\u3010\u63d2\u56fe\uff1aYOLO\u6807\u6ce8\u3011'),
    (11.3, '(iii) \u8fb9\u7f18\u521d\u7b5b', '\u9ad8\u7f6e\u4fe1\u5ea6\u672c\u5730\u5224\u5b9a',
     ORANGE, '\u3010\u63d2\u56fe\uff1a\u521d\u7b5b\u754c\u9762\u3011'),
    (14.3, '(iv) Agent\u590d\u6838', 'LangGraph ReAct\u5faa\u73af', ACCENT,
     '\u3010\u63d2\u56fe\uff1aAgent\u590d\u6838\u3011'),
    (17.3, '(v) RAG\u68c0\u7d22', 'pgvector\u8bed\u4e49\u68c0\u7d22', GREEN,
     '\u3010\u63d2\u56fe\uff1aRAG\u7ed3\u679c\u3011'),
]

for l, title, desc, color, pic_label in steps:
    add_text(l, 3.2, 2.5, 0.3, title, 14, True, DARK)
    add_box(l, 3.55, 2.5, 0.5, color, desc, 11)
    add_pic(l, 4.15, 2.5, 1.8, pic_label)

# Arrows between steps
for l in [7.9, 10.9, 13.9, 16.9]:
    add_arrow(l, 3.65, 0.3, 0.3)

# Agent calls RAG note
add_text(13.5, 6.1, 6.0, 0.3,
         '* Agent\u590d\u6838\u65f6\u8c03\u7528RAG\u68c0\u7d22\u8d28\u68c0\u6807\u51c6',
         11, False, GRAY, PP_ALIGN.CENTER)

# 4. Result feedback (bottom)
add_text(5.0, 6.6, 13.5, 0.3,
         '(vi) \u7ed3\u679c\u53cd\u9988 \u2014 SSE\u63a8\u9001\u524d\u7aef + MQTT\u4e0b\u53d1\u8fb9\u7f18',
         14, True, DARK, PP_ALIGN.CENTER)
add_box(8.0, 7.0, 7.0, 0.5, PURPLE,
        '\u7f3a\u9677\u8bb0\u5f55 + Agent\u590d\u6838\u7ed3\u8bba + \u5904\u7f6e\u5efa\u8bae', 12)
add_pic(5.0, 7.7, 13.5, 2.0,
        '\u3010\u63d2\u56fe\uff1a\u524d\u7aef\u7f3a\u9677\u8bb0\u5f55\u5217\u8868 + \u590d\u6838\u7ed3\u679c\u5c55\u793a\u3011')

# 5. Cloud Agent 5-layer architecture (left bottom)
add_text(5.0, 10.1, 8.0, 0.3,
         '\u4e91\u7aefAgent\u4e94\u5c42\u67b6\u6784', 24, True, ACCENT, PP_ALIGN.CENTER)

layers = [
    ('API\u63a5\u53e3\u5c42', 'FastAPI Routes',
     '\u8bf7\u6c42\u6821\u9a8c/\u8def\u7531\u5206\u53d1', RGBColor(0x4D, 0x4D, 0x4D)),
    ('\u670d\u52a1\u5c42', 'services/review+chat',
     '\u4e1a\u52a1\u7f16\u6392/\u9519\u8bef\u515c\u5e95', PURPLE),
    ('Agent\u7f16\u6392\u5c42', 'LangGraph ReAct',
     'Thought-Action-Observation', ACCENT),
    ('\u6a21\u578b/\u5de5\u5177/\u77e5\u8bc6', 'LLM+Tools+RAG',
     '\u70ed\u5207\u6362/\u5de5\u5177\u96c6/pgvector', BLUE),
    ('\u6570\u636e\u8bbf\u95ee\u5c42', 'SQLAlchemy+asyncpg',
     '\u8fde\u63a5\u6c60/\u5f02\u6b65ORM', GREEN),
]
ay = 10.6
for name, mod, desc, color in layers:
    add_box(5.3, ay, 3.2, 0.55, color, name, 12)
    add_text(8.7, ay + 0.02, 3.5, 0.25, mod, 11, True, DARK)
    add_text(8.7, ay + 0.28, 5.0, 0.25, desc, 10, False, GRAY)
    if ay < 12.5:
        add_arrow(6.8, ay + 0.58, 0.2, 0.25, 'down', color)
    ay += 0.65

# 6. Experimental results (right bottom)
add_text(12.5, 10.1, 6.5, 0.3, '\u5b9e\u9a8c\u7ed3\u679c', 24, True, ACCENT)

results = [
    ('\u6a21\u578b\u6027\u80fd(NEU-DET):', True),
    ('  mAP@0.5: 82.3%', False),
    ('  CPU\u63a8\u7406: ~39ms/\u5e27', False),
    ('  \u6a21\u578b: 9.5MB(FP16)', False),
    ('  \u53c2\u6570: 2.4M', False),
    ('', False),
    ('\u7aef\u5230\u7aef\u5ef6\u8fdf:', True),
    ('  \u8fb9\u7f18\u63a8\u7406: ~39ms', False),
    ('  Agent\u590d\u6838: 5-25s', False),
    ('', False),
    ('\u7cfb\u7edf\u529f\u80fd:', True),
    ('  \u2022 \u8fb9\u4e91\u534f\u540c\u521d\u7b5b\u4e0e\u590d\u6838', False),
    ('  \u2022 LLM\u8fd0\u884c\u65f6\u70ed\u5207\u6362', False),
    ('  \u2022 MQTT/HTTP\u53cc\u901a\u9053\u56de\u9000', False),
    ('  \u2022 Docker\u4e00\u952e\u90e8\u7f72', False),
]
tb = slide.shapes.add_textbox(Inches(12.5), Inches(10.6), Inches(7.0), Inches(3.8))
tf = tb.text_frame
tf.word_wrap = True
for i, (text, bold) in enumerate(results):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.name = FONT
    r.font.bold = bold
    r.font.color.rgb = DARK if bold else GRAY

# 7. Replace title/subtitle
for s in slide.shapes:
    if not s.has_text_frame:
        continue
    if s.name == 'Title 1' and abs((s.left or 0) / 914400 - 4.99) < 0.5:
        s.top = Inches(0.15)
        s.left = Inches(2.5)
        s.width = Inches(20)
        s.height = Inches(1.3)
        tf = s.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = '\u8fb9\u4e91\u534f\u540c\u667a\u80fd\u68c0\u6d4b\u7cfb\u7edf'
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = FONT
    if s.name == 'Subtitle 2':
        s.top = Inches(1.1)
        s.left = Inches(5.0)
        s.width = Inches(15)
        s.height = Inches(0.5)
        tf = s.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = 'XXX    \u8ba1\u7b97\u673a\u79d1\u5b66\u4e0e\u6280\u672f    \u6307\u5bfc\u6559\u5e08: XXX'
        r.font.size = Pt(16)
        r.font.bold = False
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        r.font.name = FONT

# 8. Left/right column text replacements
col_repl = {
    'TextBox 353': [
        '\u80cc\u666f', '',
        '\u4f20\u7edf\u5de5\u4e1a\u8d28\u68c0\u4f9d\u8d56\u4eba\u5de5\u76ee\u89c6\u5224\u65ad\uff1a',
        '  \u2022 3-5\u540d\u8d28\u68c0\u5458\u8f6e\u73ed\uff0c\u75b2\u52b3\u540e\u6f0f\u68c0\u7387\u4e0a\u534740%',
        '  \u2022 \u5224\u65ad\u6807\u51c6\u4e3b\u89c2\u4e0d\u4e00\uff0c\u7f3a\u9677\u6570\u636e\u5206\u6563',
        '  \u2022 AI\u89c6\u89c9\u6a21\u578b\u4ec5\u8f93\u51fa\u5408\u683c/\u4e0d\u5408\u683c',
        '    \u65e0\u6cd5\u8fdb\u884c\u8bed\u4e49\u63a8\u7406\u4e0e\u6839\u56e0\u5206\u6790', '',
        '\u76ee\u6807\uff1a\u6784\u5efa\u8fb9\u4e91\u534f\u540c\u7cfb\u7edf',
        '\u8fb9\u7f18\u5b9e\u65f6\u521d\u7b5b + \u4e91\u7aefAgent\u6df1\u5ea6\u590d\u6838',
    ],
    'TextBox 356': [
        '\u4f20\u7edf\u68c0\u6d4b\u65b9\u6cd5', '\u7684\u5c40\u9650\u6027', '',
        '\u4eba\u5de5\u75b2\u52b3 \u2192 \u8fde\u7eed\u5de5\u4f5c2\u5c0f\u65f6\u540e\u6f0f\u68c0\u7387\u5347\u9ad8',
        '\u6807\u51c6\u4e0d\u4e00 \u2192 \u4e0d\u540c\u8d28\u68c0\u5458\u5224\u65ad\u5b58\u5728\u4e3b\u89c2\u5dee\u5f02',
        '\u6570\u636e\u5206\u6563 \u2192 \u7f3a\u9677\u8bb0\u5f55\u65e0\u6cd5\u7cfb\u7edf\u6027\u6839\u56e0\u5206\u6790', '',
        '\u7eaf\u89c6\u89c9\u6a21\u578b\u53ea\u80fd\u4e8c\u5143\u5206\u7c7b',
        '\u65e0\u6cd5\u56de\u7b54\u201c\u7f3a\u9677\u4e25\u91cd\u7a0b\u5ea6\u201d\u201c\u4ea7\u751f\u539f\u56e0\u201d',
    ],
    'TextBox 355': [
        '\u7814\u7a76\u52a8\u673a', '',
        '\u8fb9\u4e91\u534f\u540c\u7684\u4f18\u52bf\uff1a',
        '  \u2022 \u8fb9\u7f18\u7aef\uff1a\u9ad8\u541e\u5410\u91cf\u5b9e\u65f6\u521d\u7b5b',
        '  \u2022 \u4e91\u7aef\uff1a\u4f4e\u9891\u7387\u6df1\u5ea6\u63a8\u7406',
        '  \u2022 \u5145\u5206\u5229\u7528\u4e24\u7aef\u7b97\u529b\u7279\u70b9', '',
        'LLM Agent\u63d0\u4f9b\u8bed\u4e49\u63a8\u7406\u80fd\u529b',
        '\u53ef\u7ed3\u5408RAG\u68c0\u7d22\u8d28\u68c0\u6807\u51c6\u8fdb\u884c\u590d\u6838',
    ],
    'TextBox 368': [
        '\u8fb9\u4e91\u67b6\u6784', ' \u5bf9\u6bd4', '',
        '\u8fb9\u7f18\u7aef\uff1aYOLO26n + OpenVINO',
        '  CPU\u63a8\u7406 ~39ms/\u5e27\uff0c2.4M\u53c2\u6570',
        '\u4e91\u7aef\uff1aLangChain Agent + RAG',
        '  ReAct\u5faa\u73af + pgvector\u8bed\u4e49\u68c0\u7d22',
    ],
    'TextBox 371': [
        '\u603b\u7ed3\u4e0e\u5c55\u671b', '',
        '\u672c\u8bfe\u7a0b\u8bbe\u8ba1\u5b8c\u6210\u4e86\u4e00\u5957\u8fb9\u4e91\u534f\u540c',
        '\u667a\u80fd\u68c0\u6d4b\u7cfb\u7edf\uff0c\u5b9e\u73b0\uff1a',
        '  \u2022 \u8fb9\u7f18\u7aefYOLO26n\u5b9e\u65f6\u7f3a\u9677\u68c0\u6d4b',
        '  \u2022 \u4e91\u7aefLangChain Agent\u6df1\u5ea6\u590d\u6838',
        '  \u2022 RAG\u77e5\u8bc6\u5e93\u589e\u5f3a\u8d28\u68c0\u6807\u51c6\u68c0\u7d22',
        '  \u2022 Vue 3\u53ef\u89c6\u5316\u4eba\u673a\u534f\u540c\u5e73\u53f0',
        '  \u2022 MQTT/HTTP\u53cc\u901a\u9053\u53ef\u9760\u901a\u4fe1', '',
        '\u672a\u6765\u65b9\u5411\uff1a\u591a\u4ea7\u7ebf\u5206\u5e03\u5f0f\u76d1\u63a7\u3001',
        '\u7f3a\u9677\u8d8b\u52bf\u9884\u6d4b\u6027\u7ef4\u62a4\u3001',
        '\u89c6\u89c9-\u8bed\u8a00\u591a\u6a21\u6001\u7aef\u5230\u7aef\u7f3a\u9677\u5224\u5b9a\u3002',
    ],
    'TextBox 376': ['(a)', '\u8fb9\u7f18\uff1a\u5b9e\u65f6\u521d\u7b5b'],
    'TextBox 377': ['(b)', '\u4e91\u7aef\uff1a\u6df1\u5ea6\u590d\u6838'],
}

for s in slide.shapes:
    if not s.has_text_frame:
        continue
    if s.name in col_repl:
        tf = s.text_frame
        tf.clear()
        for i, line in enumerate(col_repl[s.name]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.size = Pt(14)
            r.font.name = FONT
            r.font.bold = (i == 0)
            if i == 0:
                r.font.color.rgb = ACCENT

prs.save(dst)
print(f'Done! Saved to {dst}')
